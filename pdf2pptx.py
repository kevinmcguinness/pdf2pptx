#!/usr/bin/env python
# Copyright (c) 2020 Kevin McGuinness <kevin.mcguinness@gmail.com>

import click
import fitz
import io
import sys

from tqdm import trange
from pptx import Presentation
from pptx.util import Cm
from pathlib import Path

arg = click.argument
opt = click.option


@click.command()
@opt('-o', '--output', 'output_file', default=None,
     help='location to save the pptx (default: PDF_FILE.pptx)')
@opt('-r', '--resolution', default=300, type=int,
     help='resolution in dots per inch (default: 300)')
@opt('--from', 'start_page', default=0, type=int,
     help='first page in the pdf to copy to the pptx')
@opt('--count', 'page_count', default=None, type=int,
     help='number of pages in the pdf to copy to the pptx')
@arg('pdf_file', type=click.Path(exists=True, dir_okay=False))
def main(pdf_file, output_file, resolution, start_page, page_count):
    """
    Convert a PDF slideshow to Powerpoint PPTX.

    Renders each page as a PNG image and creates the resulting Powerpoint 
    slideshow from these images. Useful when you want to use Powerpoint 
    to present a set of PDF slides (e.g. slides from Beamer). You can then
    use the presentation capabilities of Powerpoint (notes, ink on slides,
    etc.) with slides created in LaTeX.
    """
    doc = fitz.open(pdf_file)
    print(pdf_file, 'contains', doc.pageCount, 'slides')

    if page_count is None:
        page_count = doc.pageCount

    # transformation matrix: slide to pixmap
    zoom = resolution / 72
    matrix = fitz.Matrix(zoom, zoom, 0)

    # create pptx presentation
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]

    # configure presentation aspect ratio
    page = doc.loadPage(0)
    aspect_ratio = page.rect.width / page.rect.height
    prs.slide_width = int(prs.slide_height * aspect_ratio)

    # iterate over slides
    for page_no in trange(start_page, start_page + page_count):
        page = doc.loadPage(page_no)

        # write slide as a pixmap
        pixmap = page.getPixmap(matrix=matrix)
        image_data = pixmap.getPNGData()
        image_file = io.BytesIO(image_data)
    
        # add a slide
        slide = prs.slides.add_slide(blank_slide_layout)    
        left = top = Cm(0)
        slide.shapes.add_picture(image_file, left, top, height=prs.slide_height)

    if output_file is None:
        output_file = Path(pdf_file).with_suffix('.pptx')
    
    # save presentation
    try:
        prs.save(output_file)
    except PermissionError as err:
        print(err, file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()  # pylint: disable=no-value-for-parameter
