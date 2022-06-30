# Copyright (c) 2020 Kevin McGuinness <kevin.mcguinness@gmail.com>
import fitz
import io

from tqdm import trange
from pptx import Presentation
from pptx.util import Cm
from pathlib import Path
from cli import DEFAULT_RESOLUTION, DEFAULT_START_PAGE, DEFAULT_QUIET, DEFAULT_PAGE_COUNT, DEFAULT_OUTPUT

__all__ = ['convert_pdf2pptx']


def convert_pdf2pptx(
        pdf_file, output_file=DEFAULT_OUTPUT, resolution=DEFAULT_RESOLUTION, start_page=DEFAULT_START_PAGE, page_count=DEFAULT_PAGE_COUNT,
        quiet=DEFAULT_QUIET):

    doc = fitz.open(pdf_file)
    if not quiet:
        print(pdf_file, 'contains', doc.page_count, 'slides')

    if page_count is None:
        page_count = doc.page_count

    # transformation matrix: slide to pixmap
    zoom = resolution / 72
    matrix = fitz.Matrix(zoom, zoom, 0)

    # create pptx presentation
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]

    # configure presentation aspect ratio
    page = doc.load_page(0)
    aspect_ratio = page.rect.width / page.rect.height
    prs.slide_width = int(prs.slide_height * aspect_ratio)

    # create page iterator
    if not quiet:
        page_iter = trange(start_page, start_page + page_count)
    else:
        page_iter = range(start_page, start_page + page_count)

    # iterate over slides
    for page_no in page_iter:
        page = doc.load_page(page_no)

        # write slide as a pixmap
        pixmap = page.get_pixmap(matrix=matrix)
        image_data = pixmap.tobytes(output='PNG')
        image_file = io.BytesIO(image_data)

        # add a slide
        slide = prs.slides.add_slide(blank_slide_layout)
        left = top = Cm(0)
        slide.shapes.add_picture(
            image_file, left, top, height=prs.slide_height)

    if output_file is None:
        output_file = Path(pdf_file).with_suffix('.pptx')

    # save presentation
    prs.save(output_file)
